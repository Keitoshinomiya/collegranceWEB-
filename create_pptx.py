#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""COLLEGRANCE 事業全体像 PowerPoint 生成スクリプト（16スライド版・完全最終版）"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 定数 ──
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
COLOR_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
COLOR_GREEN = RGBColor(0x3A, 0x7D, 0x44)
COLOR_RED = RGBColor(0xE6, 0x39, 0x46)
COLOR_GRAY = RGBColor(0x66, 0x66, 0x66)
COLOR_SUB = RGBColor(0x99, 0x99, 0x99)
COLOR_BODY = RGBColor(0x33, 0x33, 0x33)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
COLOR_GREEN_LIGHT = RGBColor(0xE8, 0xF5, 0xE9)
COLOR_GREEN_DARK = RGBColor(0x2E, 0x6B, 0x38)
COLOR_ORANGE = RGBColor(0xFF, 0x98, 0x00)
COLOR_BLUE = RGBColor(0x1E, 0x88, 0xE5)
COLOR_PURPLE = RGBColor(0x7B, 0x1F, 0xA2)
COLOR_TEAL = RGBColor(0x00, 0x89, 0x7B)
COLOR_MELL_GRAY = RGBColor(0x88, 0x88, 0x88)
FONT_NAME = "Meiryo"
TOTAL_SLIDES = 16

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT
blank_layout = prs.slide_layouts[6]  # blank


# ── ユーティリティ ──
def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=COLOR_BODY, alignment=PP_ALIGN.LEFT,
                font_name=FONT_NAME, line_spacing=1.2):
    """テキストボックスを追加"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    if line_spacing:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, font_size=14,
                          color=COLOR_BODY, alignment=PP_ALIGN.LEFT,
                          font_name=FONT_NAME, line_spacing=1.3):
    """複数行テキストボックス。linesは [(text, font_size, bold, color), ...] のリスト"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(item, str):
            p.text = item
            p.font.size = Pt(font_size)
            p.font.bold = False
            p.font.color.rgb = color
        else:
            t, fs, bld, clr = item
            p.text = t
            p.font.size = Pt(fs)
            p.font.bold = bld
            p.font.color.rgb = clr
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(2)
        p.line_spacing = Pt((item[1] if not isinstance(item, str) else font_size) * line_spacing)
    return txBox


def add_rect(slide, left, top, width, height, fill_color, text="",
             font_size=11, font_color=COLOR_WHITE, bold=False, alignment=PP_ALIGN.CENTER):
    """四角形を追加"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.adjustments[0] = 0.05
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = alignment
    if text:
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(font_size)
        tf.paragraphs[0].font.color.rgb = font_color
        tf.paragraphs[0].font.bold = bold
        tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    shape.text_frame.word_wrap = True
    return shape


def add_arrow(slide, left, top, width, height, color=COLOR_GREEN):
    """右向き矢印"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_down_arrow(slide, left, top, width, height, color=COLOR_GREEN):
    """下向き矢印"""
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_footer(slide):
    """フッター"""
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
                "COLLEGRANCE  |  合同会社ヤシノミ  |  Confidential",
                font_size=9, color=COLOR_SUB, alignment=PP_ALIGN.LEFT)


def add_slide_number(slide, num):
    """スライド番号"""
    add_textbox(slide, Inches(12.0), Inches(7.0), Inches(1), Inches(0.4),
                f"{num} / {TOTAL_SLIDES}", font_size=9, color=COLOR_SUB, alignment=PP_ALIGN.RIGHT)


def add_section_title(slide, title, subtitle=""):
    """スライドのセクションタイトル（上部）"""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(0.15), Inches(0.5))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_GREEN
    line.line.fill.background()
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6),
                title, font_size=26, bold=True, color=COLOR_BLACK)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.4),
                    subtitle, font_size=13, color=COLOR_GRAY)


def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None):
    """テーブルを追加。dataは[[(text, bold, color), ...], ...]"""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            d = data[r][c]
            if isinstance(d, tuple):
                text, bld, clr = d
            else:
                text, bld, clr = d, False, COLOR_BODY
            p.text = text
            p.font.size = Pt(11)
            p.font.bold = bld
            p.font.color.rgb = clr
            p.font.name = FONT_NAME
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(4)
            p.space_after = Pt(4)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_BLACK
                p.font.color.rgb = COLOR_WHITE
                p.font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_WHITE if r % 2 == 1 else COLOR_LIGHT_BG
    return table_shape


# ══════════════════════════════════════════════════════════════
# スライド1: 表紙
# ══════════════════════════════════════════════════════════════
def slide_01_cover():
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_WHITE

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_GREEN
    bar.line.fill.background()

    add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.0),
                "COLLEGRANCE", font_size=56, bold=True, color=COLOR_BLACK,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(2.9), Inches(11), Inches(0.5),
                "コレグランス", font_size=20, color=COLOR_GRAY,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
                "事業計画書  2026年4月〜2027年1月", font_size=28, bold=True, color=COLOR_GREEN,
                alignment=PP_ALIGN.CENTER)

    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(5.5), Inches(4.6), Inches(2.3), Inches(0.02))
    sep.fill.solid()
    sep.fill.fore_color.rgb = COLOR_GREEN
    sep.line.fill.background()

    add_textbox(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.4),
                "合同会社ヤシノミ", font_size=16, color=COLOR_BODY,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.4),
                "化粧品製造業・化粧品製造販売業 許可取得済み", font_size=12, color=COLOR_GRAY,
                alignment=PP_ALIGN.CENTER)

    add_footer(slide)


# ══════════════════════════════════════════════════════════════
# スライド2: 事業概要
# ══════════════════════════════════════════════════════════════
def slide_02_overview():
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "事業概要")

    add_textbox(slide, Inches(1.0), Inches(1.3), Inches(11), Inches(0.6),
                "「香りはえらべる。もっと自由に。」",
                font_size=30, bold=True, color=COLOR_GREEN, alignment=PP_ALIGN.CENTER)

    lines = [
        ("高級ブランド香水の小分け（1.5ml）お試し販売 + フルボトル直販", 16, True, COLOR_BLACK),
        ("", 10, False, COLOR_BODY),
        ("130種類以上のブランド香水を取り扱い", 14, False, COLOR_BODY),
        ("Maison Margiela / BYREDO / DIPTYQUE / DIOR / HERMES / LE LABO 他", 12, False, COLOR_GRAY),
    ]
    add_multiline_textbox(slide, Inches(1.5), Inches(2.3), Inches(10), Inches(1.8), lines,
                          alignment=PP_ALIGN.CENTER)

    # 販売チャネル
    channels = [
        ("Amazon", "小分け+ケース\n広告で集客", COLOR_BLACK),
        ("TikTok Shop", "小分け+ケース\nコンテンツ販売", COLOR_BLACK),
        ("メルカリShop", "小分け+ケース\n販売チャネル", COLOR_BLACK),
        ("自社EC (Stripe)", "フルボトル\n高利益の本命", COLOR_GREEN),
    ]
    card_w = Inches(2.5)
    gap = Inches(0.4)
    start_x = Inches(1.2)
    for i, (name, desc, bg) in enumerate(channels):
        x = start_x + i * (card_w + gap)
        y = Inches(4.0)
        shape = add_rect(slide, x, y, card_w, Inches(1.6), bg, font_size=12)
        tf = shape.text_frame
        tf.paragraphs[0].text = name
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = COLOR_WHITE
        tf.paragraphs[0].font.name = FONT_NAME
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = ""
        p2.font.size = Pt(6)
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(11)
        p3.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        p3.font.name = FONT_NAME
        p3.alignment = PP_ALIGN.CENTER
        p3.line_spacing = Pt(16)

    # 運営体制
    add_textbox(slide, Inches(1.0), Inches(5.9), Inches(11), Inches(0.4),
                "運営体制: 代表(四宮) + パート1名(みのべさん) + 11月〜(守殿さん復帰)",
                font_size=13, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)

    add_footer(slide)
    add_slide_number(slide, 2)


# ══════════════════════════════════════════════════════════════
# スライド3: ビジネスモデル全体図
# ══════════════════════════════════════════════════════════════
def slide_03_business_model():
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "ビジネスモデル全体図", "SNS認知 → Amazon購入 → LINE登録 → CRM育成 → フルボトル/クロスセル")

    # ── Row 1: メインフロー ──
    row1_y = Inches(1.6)
    box_h = Inches(1.3)
    box_w = Inches(2.2)
    arrow_w = Inches(0.5)
    arrow_h = Inches(0.3)

    # Box 1: SNS認知
    b1 = add_rect(slide, Inches(0.5), row1_y, box_w, box_h, COLOR_BLACK)
    tf = b1.text_frame
    tf.paragraphs[0].text = "SNS認知"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_WHITE
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for label in ["TikTok / Threads", "Instagram / SEO(60記事)"]:
        p = tf.add_paragraph()
        p.text = label
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

    # Instagram & Threads（点線ボックス）
    future_y = row1_y + box_h + Inches(0.15)
    for i, (label, sub) in enumerate([
        ("Instagram", "世界観構築→AI診断誘導"),
        ("Threads", "中の人投稿→認知拡大"),
    ]):
        fx = Inches(0.5) + Inches(1.15) * i
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, fx, future_y, Inches(1.1), Inches(0.65))
        shape.fill.background()
        shape.line.color.rgb = COLOR_GRAY
        shape.line.dash_style = 2
        shape.line.width = Pt(1)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = label
        tf.paragraphs[0].font.size = Pt(9)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = COLOR_GRAY
        tf.paragraphs[0].font.name = FONT_NAME
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.add_paragraph()
        p.text = sub
        p.font.size = Pt(7)
        p.font.color.rgb = COLOR_SUB
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER
    add_textbox(slide, Inches(0.5), future_y + Inches(0.68), Inches(2.3), Inches(0.25),
                "※ 今後追加予定（集客チャネル）", font_size=7, color=COLOR_SUB)

    add_arrow(slide, Inches(2.8), row1_y + Inches(0.5), arrow_w, arrow_h)

    # Box 2: Amazon小分け購入
    b2 = add_rect(slide, Inches(3.4), row1_y, Inches(2.5), box_h, COLOR_GREEN)
    tf = b2.text_frame
    tf.paragraphs[0].text = "Amazon小分け購入"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_WHITE
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for label in ["広告費 ~75万/月", "月間 ~3,500件"]:
        p = tf.add_paragraph()
        p.text = label
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0xDD, 0xFF, 0xDD)
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

    add_arrow(slide, Inches(6.0), row1_y + Inches(0.5), arrow_w, arrow_h)

    # Box 3: LINE登録
    b3 = add_rect(slide, Inches(6.6), row1_y, Inches(2.5), box_h, COLOR_GREEN_DARK)
    tf = b3.text_frame
    tf.paragraphs[0].text = "LINE登録"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_WHITE
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for label in ["ケースプレゼント施策", "登録率 5.3% / 1,300人"]:
        p = tf.add_paragraph()
        p.text = label
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0xDD, 0xFF, 0xDD)
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

    add_arrow(slide, Inches(9.2), row1_y + Inches(0.5), arrow_w, arrow_h)

    # Box 4: CRM育成
    b4 = add_rect(slide, Inches(9.8), row1_y, Inches(2.8), box_h, COLOR_BLACK)
    tf = b4.text_frame
    tf.paragraphs[0].text = "CRM育成"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_WHITE
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for label in ["LINE Harness（自社開発）", "シナリオ配信+AI診断"]:
        p = tf.add_paragraph()
        p.text = label
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

    # ── Row 2: 分岐 ──
    add_down_arrow(slide, Inches(10.3), Inches(3.1), Inches(0.3), Inches(0.5))
    add_down_arrow(slide, Inches(11.3), Inches(3.1), Inches(0.3), Inches(0.5))

    row2_y = Inches(3.8)
    b5 = add_rect(slide, Inches(9.1), row2_y, Inches(2.0), Inches(1.0), COLOR_RED)
    tf = b5.text_frame
    tf.paragraphs[0].text = "フルボトル購入"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_WHITE
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "自社EC(Stripe) / 高利益"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0xFF, 0xCC, 0xCC)
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    b6 = add_rect(slide, Inches(11.3), row2_y, Inches(1.6), Inches(1.0), COLOR_BLACK)
    tf = b6.text_frame
    tf.paragraphs[0].text = "クロスセル"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_WHITE
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "別の小分け(Amazon)"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    # ── Row 3: サステナブルループ ──
    row3_y = Inches(5.0)
    add_textbox(slide, Inches(0.5), row3_y, Inches(12), Inches(0.4),
                "サステナブルループ（ナイトライト事業）", font_size=14, bold=True, color=COLOR_GREEN)

    loop_y = Inches(5.5)
    b_src1 = add_rect(slide, Inches(0.4), loop_y, Inches(2.0), Inches(0.7), COLOR_GREEN_LIGHT,
                 text="① 小分け製造の空き瓶\n（自社ラボ発生）", font_size=9, font_color=COLOR_BLACK)
    b_src2 = add_rect(slide, Inches(0.4), loop_y + Inches(0.85), Inches(2.0), Inches(0.7), COLOR_GREEN_LIGHT,
                 text="② お客様のボトル返送\n（¥1,000引きで加工）", font_size=9, font_color=COLOR_BLACK)

    add_arrow(slide, Inches(2.5), loop_y + Inches(0.55), Inches(0.4), Inches(0.25), COLOR_GREEN)

    add_rect(slide, Inches(3.0), loop_y + Inches(0.2), Inches(2.0), Inches(0.9), COLOR_GREEN,
             text="3Dプリンタ\nナイトライト製作", font_size=11, font_color=COLOR_WHITE)

    add_arrow(slide, Inches(5.1), loop_y + Inches(0.45), Inches(0.4), Inches(0.25), COLOR_GREEN)

    add_rect(slide, Inches(5.6), loop_y, Inches(2.0), Inches(0.7), COLOR_BLACK,
             text="TikTok Shop / 自社EC\nメルカリShopで販売", font_size=9, font_color=COLOR_WHITE)
    add_rect(slide, Inches(5.6), loop_y + Inches(0.85), Inches(2.0), Inches(0.7), COLOR_BLACK,
             text="製作動画をTikTok投稿\nサステナブル×おしゃれ", font_size=9, font_color=COLOR_WHITE)

    add_arrow(slide, Inches(7.7), loop_y + Inches(0.55), Inches(0.4), Inches(0.25), COLOR_GREEN)

    add_rect(slide, Inches(8.2), loop_y + Inches(0.05), Inches(2.2), Inches(0.6), COLOR_GREEN_LIGHT,
             text="新規認知→小分け購入→ループ", font_size=10, font_color=COLOR_GREEN_DARK)

    # 返送→フルボトル購入動機の循環
    add_rect(slide, Inches(8.2), loop_y + Inches(0.75), Inches(2.2), Inches(0.7), COLOR_ORANGE,
             text="ボトル返送→¥1,000引き加工\n→フルボトル購入動機", font_size=9, font_color=COLOR_WHITE, bold=True)

    add_down_arrow(slide, Inches(9.8), Inches(4.9), Inches(0.3), Inches(0.3))

    add_footer(slide)
    add_slide_number(slide, 3)


# ══════════════════════════════════════════════════════════════
# スライド4: 競合分析
# ══════════════════════════════════════════════════════════════
def slide_04_competitor():
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "競合分析", "主要競合: MELL fragrance（メルフレグランス）")

    # 比較テーブル
    comp_data = [
        [("項目", True, COLOR_WHITE), ("MELL fragrance", True, COLOR_WHITE),
         ("COLLEGRANCE", True, COLOR_WHITE), ("差分", True, COLOR_WHITE)],
        [("メルカリ累計", False, COLOR_BODY),
         ("17,374件", False, COLOR_MELL_GRAY),
         ("ほぼゼロ", False, COLOR_RED),
         ("要開拓", True, COLOR_RED)],
        [("メルカリ日販", False, COLOR_BODY),
         ("40〜50件/日", False, COLOR_MELL_GRAY),
         ("-", False, COLOR_BODY),
         ("", False, COLOR_BODY)],
        [("TikTok Shop", False, COLOR_BODY),
         ("~10,000件(5ヶ月)", False, COLOR_MELL_GRAY),
         ("25件(4月ローンチ)", False, COLOR_RED),
         ("要成長", True, COLOR_RED)],
        [("TikTokレビュー", False, COLOR_BODY),
         ("多数", False, COLOR_MELL_GRAY),
         ("1件", False, COLOR_RED),
         ("要50件", True, COLOR_RED)],
        [("Amazon", False, COLOR_BODY),
         ("後発", False, COLOR_MELL_GRAY),
         ("1位", True, COLOR_GREEN),
         ("★優位", True, COLOR_GREEN)],
        [("品揃え", False, COLOR_BODY),
         ("多い", False, COLOR_MELL_GRAY),
         ("130種類+", False, COLOR_BODY),
         ("", False, COLOR_BODY)],
        [("許可", False, COLOR_BODY),
         ("不明", False, COLOR_MELL_GRAY),
         ("★製造販売業許可", True, COLOR_GREEN),
         ("差別化ポイント", True, COLOR_GREEN)],
        [("差別化", False, COLOR_BODY),
         ("品数", False, COLOR_MELL_GRAY),
         ("ケース・LINE CRM・AI診断・ナイトライト", False, COLOR_GREEN),
         ("", False, COLOR_BODY)],
    ]

    add_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(4.2),
              9, 4, comp_data,
              col_widths=[Inches(2.0), Inches(3.0), Inches(4.5), Inches(2.8)])

    # COLLEGRANCEの強み（右下エリア）
    rx = Inches(7.0)
    sy = Inches(5.6)
    add_textbox(slide, rx, sy, Inches(5.5), Inches(0.3),
                "COLLEGRANCEの強み", font_size=14, bold=True, color=COLOR_GREEN)
    strengths = [
        "1. Amazon1位（デザイン性・切り口で後発ながら獲得）",
        "2. 化粧品製造販売業許可（信頼性）",
        "3. LINE CRM（リピート・アップセル仕組み）",
        "4. AI香水診断（自社サイト独自機能）",
        "5. ナイトライト（サステナブル×ものづくり）",
    ]
    for i, s in enumerate(strengths):
        add_textbox(slide, rx + Inches(0.2), sy + Inches(0.35) + i * Inches(0.25), Inches(5.3), Inches(0.25),
                    s, font_size=10, color=COLOR_GREEN_DARK)

    # 課題
    add_rect(slide, Inches(0.5), Inches(5.6), Inches(6.0), Inches(1.2), COLOR_RED,
             font_size=10, font_color=COLOR_WHITE)
    shape = slide.shapes[-1]
    tf = shape.text_frame
    tf.paragraphs[0].text = "課題: TikTok・メルカリのレビュー数で大きく負けている"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_WHITE
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "→ レビュー獲得施策が最優先"
    p2.font.size = Pt(13)
    p2.font.color.rgb = RGBColor(0xFF, 0xCC, 0xCC)
    p2.font.name = FONT_NAME
    p2.alignment = PP_ALIGN.CENTER

    add_footer(slide)
    add_slide_number(slide, 4)


# ══════════════════════════════════════════════════════════════
# スライド5: Amazon販売実績
# ══════════════════════════════════════════════════════════════
def slide_05_amazon_sales():
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Amazon販売実績", "2025年10月〜2026年3月（実測+推定値）")

    table_data = [
        [("月", True, COLOR_WHITE), ("注文数", True, COLOR_WHITE),
         ("売上(税込)", True, COLOR_WHITE), ("広告費(税抜)", True, COLOR_WHITE),
         ("ROAS", True, COLOR_WHITE)],
        [("10月", True, COLOR_BODY), ("1,587", False, COLOR_BODY),
         ("¥1,712,373", False, COLOR_BODY), ("¥829,402", False, COLOR_BODY),
         ("1.70", False, COLOR_RED)],
        [("11月", True, COLOR_BODY), ("2,494", False, COLOR_BODY),
         ("¥2,691,026", False, COLOR_BODY), ("¥758,425", False, COLOR_BODY),
         ("2.31", False, COLOR_BODY)],
        [("12月", True, COLOR_BODY), ("4,941", False, COLOR_BODY),
         ("¥5,331,339", False, COLOR_GREEN), ("¥802,605", False, COLOR_BODY),
         ("4.18", True, COLOR_GREEN)],
        [("1月", True, COLOR_BODY), ("3,433", False, COLOR_BODY),
         ("¥3,704,207", False, COLOR_BODY), ("¥618,290", False, COLOR_BODY),
         ("3.83", False, COLOR_BODY)],
        [("2月", True, COLOR_BODY), ("4,378", False, COLOR_BODY),
         ("¥4,723,862", False, COLOR_BODY), ("¥839,654", False, COLOR_BODY),
         ("4.01", False, COLOR_GREEN)],
        [("3月", True, COLOR_BODY), ("4,323", False, COLOR_BODY),
         ("¥4,664,517", False, COLOR_BODY), ("¥702,310", False, COLOR_BODY),
         ("3.58", False, COLOR_BODY)],
        [("合計", True, COLOR_BLACK), ("21,156", True, COLOR_BLACK),
         ("¥22,827,324", True, COLOR_BLACK), ("¥4,550,686", True, COLOR_BLACK),
         ("3.36", True, COLOR_GREEN)],
    ]

    add_table(slide, Inches(0.5), Inches(1.4), Inches(7.5), Inches(4.0),
              8, 5, table_data,
              col_widths=[Inches(1.0), Inches(1.2), Inches(2.0), Inches(1.8), Inches(1.0)])

    # 棒グラフ
    months_bar = [
        ("10月", 1587), ("11月", 2494), ("12月", 4941),
        ("1月", 3433), ("2月", 4378), ("3月", 4323),
    ]
    max_orders = 4941
    bar_max_h = Inches(2.8)
    bar_w = Inches(0.65)
    gap = Inches(0.15)
    start_x = Inches(8.5)
    base_y = Inches(5.2)

    add_textbox(slide, Inches(8.5), Inches(1.4), Inches(4), Inches(0.3),
                "月別注文数推移", font_size=14, bold=True, color=COLOR_GREEN)

    for i, (month, orders) in enumerate(months_bar):
        x = start_x + i * (bar_w + gap)
        ratio = orders / max_orders
        bar_h = bar_max_h * ratio
        bar_y = base_y - bar_h
        is_peak = (month == "12月")
        fill = COLOR_RED if is_peak else COLOR_GREEN
        add_rect(slide, x, bar_y, bar_w, bar_h, fill)
        add_textbox(slide, x - Inches(0.1), bar_y - Inches(0.35), bar_w + Inches(0.2), Inches(0.3),
                    f"{orders:,}", font_size=9, bold=True,
                    color=COLOR_RED if is_peak else COLOR_BLACK,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x - Inches(0.1), base_y + Inches(0.05), bar_w + Inches(0.2), Inches(0.3),
                    month, font_size=9, bold=True, color=COLOR_BLACK,
                    alignment=PP_ALIGN.CENTER)

    # サマリー + TACOS追加
    lines = [
        ("ROAS推移", 13, True, COLOR_GRAY),
        ("10月: 1.70 → 3月: 3.58", 16, True, COLOR_GREEN),
        ("（広告効率 2.1倍に改善）", 11, False, COLOR_GRAY),
        ("", 8, False, COLOR_BODY),
        ("TACOS: 21.9% → 目標15%", 14, True, COLOR_RED),
        ("", 8, False, COLOR_BODY),
        ("平均単価（税込）: ¥1,079", 12, False, COLOR_BLACK),
    ]
    add_multiline_textbox(slide, Inches(8.5), Inches(5.5), Inches(4), Inches(1.5), lines)

    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(7.5), Inches(0.6),
                "※10月・11月は実測値、12月〜3月はアンケート回答率からの推定\n※広告データは全月セラーセントラル実測値（税抜）、売上はSP-API（税込）",
                font_size=9, color=COLOR_SUB)

    add_footer(slide)
    add_slide_number(slide, 5)


# ══════════════════════════════════════════════════════════════
# スライド6: ユニットエコノミクス
# ══════════════════════════════════════════════════════════════
def slide_06_unit_economics():
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "ユニットエコノミクス", "1個あたりの利益構造")

    # ── 左: Amazon小分け ──
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(3.0), Inches(0.3),
                "Amazon小分け（1.5ml）", font_size=16, bold=True, color=COLOR_GREEN)

    econ_items = [
        ("販売価格(税抜)", "¥981", COLOR_BLACK, True),
        ("- 香水原価+副資材", "¥230", COLOR_GRAY, False),
        ("- FBA出荷費用", "¥202", COLOR_GRAY, False),
        ("- Amazon手数料(8%)", "¥78", COLOR_GRAY, False),
        ("= 利益（広告費前）", "¥470", COLOR_GREEN, True),
    ]
    for i, (label, value, color, bold) in enumerate(econ_items):
        y = Inches(1.8) + i * Inches(0.4)
        if label.startswith("="):
            sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         Inches(0.7), y - Inches(0.05), Inches(2.8), Inches(0.02))
            sep.fill.solid()
            sep.fill.fore_color.rgb = COLOR_GREEN
            sep.line.fill.background()
        add_textbox(slide, Inches(0.7), y, Inches(2.0), Inches(0.3),
                    label, font_size=11, bold=bold, color=color)
        add_textbox(slide, Inches(2.7), y, Inches(0.8), Inches(0.3),
                    value, font_size=11, bold=bold, color=color, alignment=PP_ALIGN.RIGHT)

    add_rect(slide, Inches(0.5), Inches(4.0), Inches(3.2), Inches(0.8), COLOR_GREEN_LIGHT,
             font_size=10, font_color=COLOR_BLACK)
    shape = slide.shapes[-1]
    tf = shape.text_frame
    tf.paragraphs[0].text = "ケース付き(30%): +¥280 → 利益¥750/個"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_GREEN_DARK
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "加重平均利益: ¥554/個"
    p2.font.size = Pt(12)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_BLACK
    p2.font.name = FONT_NAME
    p2.alignment = PP_ALIGN.CENTER

    # ── 中央: フルボトル ──
    mx = Inches(4.0)
    add_textbox(slide, mx, Inches(1.4), Inches(3.0), Inches(0.3),
                "フルボトル（自社EC）", font_size=16, bold=True, color=COLOR_RED)

    fb_items = [
        ("価格設定", "仕入原価 x 1.25 x 1.10"),
        ("粗利率", "25%"),
        ("在庫リスク", "ゼロ（受注後発注）"),
        ("広告費", "¥0"),
        ("Amazon手数料", "¥0"),
    ]
    for i, (label, value) in enumerate(fb_items):
        y = Inches(1.8) + i * Inches(0.42)
        add_textbox(slide, mx, y, Inches(1.3), Inches(0.3),
                    label, font_size=10, bold=True, color=COLOR_GRAY)
        add_textbox(slide, mx + Inches(1.3), y, Inches(2.0), Inches(0.3),
                    value, font_size=10, color=COLOR_RED if label == "粗利率" else COLOR_BODY)

    add_rect(slide, mx, Inches(4.0), Inches(3.0), Inches(0.8), COLOR_RED,
             text="低マージンだが在庫リスクゼロ\nLINE経由でAmazon広告費不要", font_size=10,
             font_color=COLOR_WHITE, bold=True)

    # ── 右: ナイトライト2パターン ──
    rx = Inches(7.3)
    add_textbox(slide, rx, Inches(1.4), Inches(5.5), Inches(0.3),
                "ナイトライト（新規）", font_size=16, bold=True, color=COLOR_ORANGE)

    # パターン1: 自社ボトル
    add_rect(slide, rx, Inches(1.8), Inches(2.7), Inches(1.8), COLOR_GREEN_LIGHT,
             font_size=9, font_color=COLOR_BLACK)
    shape = slide.shapes[-1]
    tf = shape.text_frame
    tf.paragraphs[0].text = "自社ボトル"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_BLACK
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for line_text in ["¥3,480 - ¥300 = 粗利¥3,180", "粗利率 91%"]:
        p = tf.add_paragraph()
        p.text = line_text
        p.font.size = Pt(10)
        p.font.bold = "91%" in line_text
        p.font.color.rgb = COLOR_ORANGE if "91%" in line_text else COLOR_BODY
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

    # パターン2: 返送加工
    add_rect(slide, rx + Inches(2.9), Inches(1.8), Inches(2.7), Inches(1.8), COLOR_GREEN_LIGHT,
             font_size=9, font_color=COLOR_BLACK)
    shape = slide.shapes[-1]
    tf = shape.text_frame
    tf.paragraphs[0].text = "返送加工"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_BLACK
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for line_text in ["¥2,480 - ¥300 = 粗利¥2,180", "粗利率 88%", "※ボトル返送者は¥1,000引き"]:
        p = tf.add_paragraph()
        p.text = line_text
        p.font.size = Pt(10) if "※" not in line_text else Pt(8)
        p.font.bold = "88%" in line_text
        p.font.color.rgb = COLOR_ORANGE if "88%" in line_text else (COLOR_GRAY if "※" in line_text else COLOR_BODY)
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

    add_rect(slide, rx, Inches(4.0), Inches(5.6), Inches(0.8), COLOR_ORANGE,
             text="超高利益率 + ボトル返送→フルボトル購入動機に繋がる循環", font_size=12,
             font_color=COLOR_WHITE, bold=True)

    # 下部サマリー
    add_rect(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.8), COLOR_GREEN_LIGHT,
             font_size=10, font_color=COLOR_BLACK)
    shape = slide.shapes[-1]
    tf = shape.text_frame
    tf.paragraphs[0].text = "Amazon小分け = キャッシュフローエンジン  /  フルボトル = ノーリスク追加利益  /  ナイトライト = 高利益の成長ドライバー"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_GREEN_DARK
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_footer(slide)
    add_slide_number(slide, 6)


# ══════════════════════════════════════════════════════════════
# スライド7: 月別P/L（実績）
# ══════════════════════════════════════════════════════════════
def slide_07_monthly_pl():
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "Amazon小分け事業 月別P/L", "2025年10月〜2026年3月")

    pl_data = [
        [("月", True, COLOR_WHITE), ("売上(税抜)", True, COLOR_WHITE),
         ("原価+手数料", True, COLOR_WHITE), ("広告費", True, COLOR_WHITE),
         ("粗利益", True, COLOR_WHITE), ("利益率", True, COLOR_WHITE)],
        [("10月", True, COLOR_BODY), ("¥1,556,703", False, COLOR_BODY),
         ("¥810,215", False, COLOR_BODY), ("¥829,402", False, COLOR_BODY),
         ("-¥82,915", True, COLOR_RED), ("-5.3%", True, COLOR_RED)],
        [("11月", True, COLOR_BODY), ("¥2,446,387", False, COLOR_BODY),
         ("¥1,273,269", False, COLOR_BODY), ("¥758,425", False, COLOR_BODY),
         ("¥414,694", True, COLOR_GREEN), ("17.0%", True, COLOR_GREEN)],
        [("12月", True, COLOR_BODY), ("¥4,846,672", False, COLOR_BODY),
         ("¥2,522,542", False, COLOR_BODY), ("¥802,605", False, COLOR_BODY),
         ("¥1,521,525", True, COLOR_GREEN), ("31.4%", True, COLOR_GREEN)],
        [("1月", True, COLOR_BODY), ("¥3,367,461", False, COLOR_BODY),
         ("¥1,752,659", False, COLOR_BODY), ("¥618,290", False, COLOR_BODY),
         ("¥996,512", True, COLOR_GREEN), ("29.6%", True, COLOR_GREEN)],
        [("2月", True, COLOR_BODY), ("¥4,294,420", False, COLOR_BODY),
         ("¥2,235,113", False, COLOR_BODY), ("¥839,654", False, COLOR_BODY),
         ("¥1,219,654", True, COLOR_GREEN), ("28.4%", True, COLOR_GREEN)],
        [("3月", True, COLOR_BODY), ("¥4,240,470", False, COLOR_BODY),
         ("¥2,207,033", False, COLOR_BODY), ("¥702,310", False, COLOR_BODY),
         ("¥1,331,127", True, COLOR_GREEN), ("31.4%", True, COLOR_GREEN)],
        [("合計", True, COLOR_BLACK), ("¥20,752,113", True, COLOR_BLACK),
         ("¥10,800,831", True, COLOR_BLACK), ("¥4,550,686", True, COLOR_BLACK),
         ("¥5,400,596", True, COLOR_GREEN), ("26.0%", True, COLOR_GREEN)],
    ]

    add_table(slide, Inches(0.5), Inches(1.4), Inches(11.5), Inches(4.5),
              8, 6, pl_data,
              col_widths=[Inches(1.0), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.5), Inches(1.5)])

    cards = [
        ("月平均利益", "¥900,099", COLOR_GREEN),
        ("人件費控除後", "¥740,099/月", COLOR_BLACK),
        ("安定黒字化", "11月以降", COLOR_GREEN_DARK),
    ]
    for i, (label, value, color) in enumerate(cards):
        x = Inches(0.5) + i * Inches(4.0)
        shape = add_rect(slide, x, Inches(6.0), Inches(3.5), Inches(0.7), color)
        tf = shape.text_frame
        tf.paragraphs[0].text = label
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = False
        tf.paragraphs[0].font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        tf.paragraphs[0].font.name = FONT_NAME
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = value
        p2.font.size = Pt(20)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE
        p2.font.name = FONT_NAME
        p2.alignment = PP_ALIGN.CENTER

    add_textbox(slide, Inches(0.5), Inches(6.8), Inches(11), Inches(0.3),
                "※10月は立ち上げ月で赤字。パート人件費 月16万円。12月〜3月の注文数はアンケート回答率からの推定値",
                font_size=9, color=COLOR_SUB)

    add_footer(slide)
    add_slide_number(slide, 7)


# ══════════════════════════════════════════════════════════════
# スライド8: アンケートファネル分析
# ══════════════════════════════════════════════════════════════
def slide_08_funnel():
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "アンケートファネル分析", "Amazon購入者 → LINE登録 → クロスセル/アップセル対象")

    funnel_data = [
        ("Amazon購入者（推定）", "21,156人", "100%", Inches(10.0), COLOR_BLACK),
        ("LINE登録（アンケート回答）", "1,258人", "登録率 5.3%", Inches(8.0), COLOR_GREEN),
        ("「他にも試したい」", "1,087人", "86.4%", Inches(6.5), COLOR_GREEN_DARK),
        ("「フルボトル検討中」", "728人", "57.9%", Inches(5.0), COLOR_RED),
    ]

    center_x = Inches(4.5)
    y_start = Inches(1.5)
    step_h = Inches(1.2)

    for i, (label, count, rate, fw, color) in enumerate(funnel_data):
        y = y_start + i * step_h
        x = center_x - fw / 2
        shape = add_rect(slide, x, y, fw, Inches(0.9), color)
        tf = shape.text_frame
        tf.paragraphs[0].text = f"{label}    {count}    ({rate})"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = COLOR_WHITE
        tf.paragraphs[0].font.name = FONT_NAME
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        if i < len(funnel_data) - 1:
            add_down_arrow(slide, center_x - Inches(0.15), y + Inches(0.9),
                           Inches(0.3), Inches(0.25), COLOR_GRAY)

    rx = Inches(10.0)
    add_rect(slide, rx, Inches(3.9), Inches(2.8), Inches(1.0), COLOR_GREEN_DARK,
             text="クロスセル対象\n1,087人", font_size=13, font_color=COLOR_WHITE, bold=True)
    add_rect(slide, rx, Inches(5.1), Inches(2.8), Inches(1.0), COLOR_RED,
             text="アップセル対象\n728人", font_size=13, font_color=COLOR_WHITE, bold=True)

    # LINE配信ローンチ追記
    lines = [
        ("★ 728人がフルボトル購入を検討中", 16, True, COLOR_RED),
        ("→ この2セグメントに対するLINE配信を明日ローンチ", 14, True, COLOR_GREEN),
        ("→ フルボトルLP / クロスセルLPをパーソナライズ配信", 12, False, COLOR_BODY),
    ]
    add_multiline_textbox(slide, Inches(1.0), Inches(6.0), Inches(10), Inches(1.0), lines,
                          alignment=PP_ALIGN.LEFT)

    add_footer(slide)
    add_slide_number(slide, 8)


# ══════════════════════════════════════════════════════════════
# スライド9: 顧客属性
# ══════════════════════════════════════════════════════════════
def slide_09_customer():
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "顧客属性（アンケートデータ）", "回答者: 1,258人")

    lx = Inches(0.8)

    # 性別
    add_textbox(slide, lx, Inches(1.5), Inches(3), Inches(0.3),
                "性別", font_size=14, bold=True, color=COLOR_GREEN)
    gender_data = [("男性 70%", 0.70, COLOR_BLACK), ("女性 29%", 0.29, COLOR_GREEN), ("その他 1%", 0.01, COLOR_GRAY)]
    bar_base_x = lx
    bar_y = Inches(1.9)
    total_w = Inches(5.0)
    cum_x = bar_base_x
    for label, ratio, color in gender_data:
        w = int(total_w * ratio)
        add_rect(slide, cum_x, bar_y, w, Inches(0.45), color, text=label,
                 font_size=10, font_color=COLOR_WHITE, bold=True)
        cum_x += w

    # 年齢
    add_textbox(slide, lx, Inches(2.6), Inches(3), Inches(0.3),
                "年齢分布", font_size=14, bold=True, color=COLOR_GREEN)
    ages = [("50代", 25), ("40代", 20), ("30代", 19), ("20代", 18), ("60代〜", 12), ("10代", 6)]
    age_bar_y = Inches(3.0)
    age_bar_max_w = Inches(4.0)
    for i, (age, pct) in enumerate(ages):
        y = age_bar_y + i * Inches(0.38)
        add_textbox(slide, lx, y, Inches(0.8), Inches(0.3), age,
                    font_size=10, bold=True, color=COLOR_BODY, alignment=PP_ALIGN.RIGHT)
        w = int(age_bar_max_w * pct / 25)
        add_rect(slide, lx + Inches(0.9), y + Inches(0.03), w, Inches(0.3),
                 COLOR_GREEN, text=f"{pct}%", font_size=9, font_color=COLOR_WHITE)

    # 満足度
    mx = Inches(6.5)
    add_textbox(slide, mx, Inches(1.5), Inches(3), Inches(0.3),
                "満足度", font_size=14, bold=True, color=COLOR_GREEN)

    add_textbox(slide, mx, Inches(1.9), Inches(3), Inches(0.5),
                "★4以上: 88.9%", font_size=28, bold=True, color=COLOR_RED)

    stars = [("★5", 59.4, COLOR_RED), ("★4", 29.5, COLOR_GREEN), ("★3", 9.5, COLOR_GRAY),
             ("★2", 1.1, COLOR_GRAY), ("★1", 0.5, COLOR_GRAY)]
    star_y = Inches(2.7)
    for i, (label, pct, color) in enumerate(stars):
        y = star_y + i * Inches(0.35)
        add_textbox(slide, mx, y, Inches(0.5), Inches(0.3), label,
                    font_size=10, bold=True, color=COLOR_BODY)
        w = int(Inches(3.5) * pct / 60)
        add_rect(slide, mx + Inches(0.6), y + Inches(0.03), w, Inches(0.25),
                 color, text=f"{pct}%", font_size=8, font_color=COLOR_WHITE)

    # 香水を知ったきっかけ
    add_textbox(slide, mx, Inches(4.8), Inches(5), Inches(0.3),
                "香水を知ったきっかけ", font_size=14, bold=True, color=COLOR_GREEN)
    sources = [("YouTube", 29), ("Amazon・楽天", 26), ("SNS(IG/TikTok)", 22),
               ("友人・知人", 12), ("雑誌", 6), ("その他", 5)]
    src_y = Inches(5.2)
    for i, (src, pct) in enumerate(sources):
        y = src_y + i * Inches(0.32)
        add_textbox(slide, mx, y, Inches(1.8), Inches(0.3), src,
                    font_size=10, color=COLOR_BODY, alignment=PP_ALIGN.RIGHT)
        w = int(Inches(3.0) * pct / 30)
        add_rect(slide, mx + Inches(2.0), y + Inches(0.03), w, Inches(0.25),
                 COLOR_BLACK, text=f"{pct}%", font_size=8, font_color=COLOR_WHITE)

    # 人気商品
    add_textbox(slide, lx, Inches(5.3), Inches(5), Inches(0.3),
                "人気商品TOP3", font_size=14, bold=True, color=COLOR_GREEN)
    prods = [
        ("1位", "商品030", "42.1%"),
        ("2位", "商品018", "20.6%"),
        ("3位", "商品039", "7.8%"),
    ]
    for i, (rank, name, pct) in enumerate(prods):
        y = Inches(5.7) + i * Inches(0.35)
        add_textbox(slide, lx, y, Inches(0.6), Inches(0.3), rank,
                    font_size=12, bold=True, color=COLOR_RED)
        add_textbox(slide, lx + Inches(0.6), y, Inches(1.5), Inches(0.3), name,
                    font_size=12, color=COLOR_BODY)
        add_textbox(slide, lx + Inches(2.2), y, Inches(1.0), Inches(0.3), pct,
                    font_size=12, bold=True, color=COLOR_BLACK)

    add_footer(slide)
    add_slide_number(slide, 9)


# ══════════════════════════════════════════════════════════════
# スライド10: 販売チャネル詳細
# ══════════════════════════════════════════════════════════════
def slide_10_channels():
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "販売チャネル詳細")

    data = [
        [("チャネル", True, COLOR_WHITE), ("商品", True, COLOR_WHITE),
         ("現状", True, COLOR_WHITE), ("発送", True, COLOR_WHITE),
         ("KPI", True, COLOR_WHITE)],
        [("Amazon", True, COLOR_BODY), ("小分け+ケース", False, COLOR_BODY),
         ("月3,500件/1位", True, COLOR_GREEN), ("FBA", False, COLOR_BODY),
         ("TACOS 15%", False, COLOR_BODY)],
        [("TikTok Shop", True, COLOR_BODY), ("小分け+ケース", False, COLOR_BODY),
         ("25件/月", False, COLOR_BODY), ("ネコポス", False, COLOR_BODY),
         ("SPS30件→レビュー50件", False, COLOR_BODY)],
        [("メルカリShop", True, COLOR_BODY), ("小分け+ケース", False, COLOR_BODY),
         ("8SKU/立ち上げ中", False, COLOR_BODY), ("ネコポス", False, COLOR_BODY),
         ("レビュー獲得→日販20件", False, COLOR_BODY)],
        [("自社EC", True, COLOR_RED), ("フルボトル", False, COLOR_BODY),
         ("5件累計", False, COLOR_BODY), ("ヤマト運輸", False, COLOR_BODY),
         ("LINE転換率20%", False, COLOR_RED)],
        [("TikTok/自社", True, COLOR_ORANGE), ("ナイトライト", False, COLOR_BODY),
         ("準備中", False, COLOR_GRAY), ("ヤマト運輸", False, COLOR_BODY),
         ("月100個", False, COLOR_ORANGE)],
    ]

    add_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(3.5),
              6, 5, data,
              col_widths=[Inches(2.0), Inches(2.0), Inches(2.8), Inches(2.0), Inches(3.5)])

    # 運用体制
    add_textbox(slide, Inches(0.5), Inches(5.0), Inches(12), Inches(0.3),
                "運用・オペレーション体制", font_size=16, bold=True, color=COLOR_GREEN)

    ops_items = [
        ("発送管理", "ネクストエンジン（TikTok+メルカリ統合）", COLOR_BODY),
        ("オペレーション", "固まるまで四宮が対応", COLOR_BODY),
        ("11月〜", "守殿さん復帰で体制強化", COLOR_GREEN),
    ]
    for i, (label, desc, color) in enumerate(ops_items):
        y = Inches(5.4) + i * Inches(0.4)
        add_rect(slide, Inches(0.5), y, Inches(0.15), Inches(0.25), COLOR_GREEN)
        add_textbox(slide, Inches(0.8), y, Inches(1.5), Inches(0.3),
                    label, font_size=12, bold=True, color=COLOR_BLACK)
        add_textbox(slide, Inches(2.5), y, Inches(8), Inches(0.3),
                    desc, font_size=12, color=color)

    add_footer(slide)
    add_slide_number(slide, 10)


# ══════════════════════════════════════════════════════════════
# スライド11: LINE CRM設計
# ══════════════════════════════════════════════════════════════
def slide_11_line_crm():
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "LINE CRM設計", "LINE Harness（自社開発CRM/MA）")

    # 明日ローンチ強調
    add_rect(slide, Inches(7.0), Inches(0.3), Inches(3.0), Inches(0.5), COLOR_RED,
             text="明日ローンチ!", font_size=16, font_color=COLOR_WHITE, bold=True)

    scenario = [
        ("Day 0", "ウェルカム + 購入商品タグ付け", COLOR_GREEN),
        ("Day 3", "「香りはいかがですか？」", COLOR_GREEN_DARK),
        ("Day 7", "フルボトルLP配信\ncollegrance.com/?from=sample&product={ID}", COLOR_RED),
        ("Day 14", "AI香水診断誘導\ncollegrance.com/", COLOR_BLACK),
    ]

    add_textbox(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.3),
                "シナリオ配信", font_size=16, bold=True, color=COLOR_GREEN)

    for i, (day, desc, color) in enumerate(scenario):
        y = Inches(1.9) + i * Inches(1.1)
        add_rect(slide, Inches(0.8), y, Inches(1.0), Inches(0.5), color,
                 text=day, font_size=12, font_color=COLOR_WHITE, bold=True)
        add_textbox(slide, Inches(2.0), y, Inches(4.5), Inches(0.8),
                    desc, font_size=12, color=COLOR_BODY)
        if i < len(scenario) - 1:
            add_down_arrow(slide, Inches(1.15), y + Inches(0.55),
                           Inches(0.2), Inches(0.3), COLOR_GRAY)

    rx = Inches(7.0)
    add_textbox(slide, rx, Inches(1.4), Inches(5), Inches(0.3),
                "アンケート → レコメンド", font_size=16, bold=True, color=COLOR_GREEN)

    reco_flow = [
        ("アンケート", "「次に試したい香りは？」", COLOR_BLACK),
        ("タグ付け", "回答に基づき興味タグ付与", COLOR_GREEN_DARK),
        ("レコメンド配信", "小分けおすすめLP\n?from=recommend&product={ID}", COLOR_GREEN),
    ]
    for i, (title, desc, color) in enumerate(reco_flow):
        y = Inches(1.9) + i * Inches(1.3)
        add_rect(slide, rx, y, Inches(2.0), Inches(0.5), color,
                 text=title, font_size=11, font_color=COLOR_WHITE, bold=True)
        add_textbox(slide, rx + Inches(2.2), y, Inches(3.5), Inches(0.8),
                    desc, font_size=11, color=COLOR_BODY)
        if i < len(reco_flow) - 1:
            add_down_arrow(slide, rx + Inches(0.85), y + Inches(0.55),
                           Inches(0.2), Inches(0.3), COLOR_GRAY)

    add_rect(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.8), COLOR_GREEN_LIGHT,
             text="ポイント: 商品IDに応じてURLを動的生成し、パーソナライズ配信を実現",
             font_size=14, font_color=COLOR_BLACK, bold=True)

    add_footer(slide)
    add_slide_number(slide, 11)


# ══════════════════════════════════════════════════════════════
# スライド12: SNS + 自社サイト
# ══════════════════════════════════════════════════════════════
def slide_12_site_sns():
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "自社サイト + SNS戦略")

    add_textbox(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.3),
                "collegrance.com（リニューアル完了）", font_size=16, bold=True, color=COLOR_GREEN)

    site_features = [
        "AI香水診断（Claude AI）",
        "フルボトルEC（130種類、Stripe決済）",
        "ブログ60本（SEO、100本目標）",
        "商品LP（LINE配信のリンク先）",
        "GA4 eコマーストラッキング",
    ]
    for i, feat in enumerate(site_features):
        y = Inches(1.9) + i * Inches(0.4)
        add_rect(slide, Inches(0.8), y, Inches(0.15), Inches(0.25), COLOR_GREEN)
        add_textbox(slide, Inches(1.1), y, Inches(5), Inches(0.3),
                    feat, font_size=13, color=COLOR_BODY)

    # SNS戦略
    add_textbox(slide, Inches(6.5), Inches(1.4), Inches(6), Inches(0.3),
                "SNS戦略", font_size=16, bold=True, color=COLOR_GREEN)

    sns_data = [
        [("チャネル", True, COLOR_WHITE), ("担当", True, COLOR_WHITE),
         ("状態", True, COLOR_WHITE)],
        [("TikTok", True, COLOR_BODY), ("みのべさん", False, COLOR_BODY),
         ("稼働中（毎日1本）", True, COLOR_GREEN)],
        [("Threads", True, COLOR_BODY), ("keito（中の人）", False, COLOR_BODY),
         ("準備中", False, COLOR_GRAY)],
        [("Instagram", True, COLOR_BODY), ("keito", False, COLOR_BODY),
         ("準備中", False, COLOR_GRAY)],
    ]
    add_table(slide, Inches(6.5), Inches(1.9), Inches(6.0), Inches(2.0),
              4, 3, sns_data,
              col_widths=[Inches(1.5), Inches(2.0), Inches(2.5)])

    # TikTok動画の現状追加
    add_rect(slide, Inches(6.5), Inches(4.1), Inches(6.0), Inches(0.8), COLOR_GREEN_LIGHT,
             font_size=10, font_color=COLOR_BLACK)
    shape = slide.shapes[-1]
    tf = shape.text_frame
    tf.paragraphs[0].text = "TikTok: みのべさん毎日1本投稿"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_BLACK
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    p2 = tf.add_paragraph()
    p2.text = "現在の平均再生数: 約2,000回/本 / アルゴリズム研究中"
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_GRAY
    p2.font.name = FONT_NAME
    p2.alignment = PP_ALIGN.LEFT

    # ナイトライト事業
    add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.3),
                "新施策: ナイトライト事業", font_size=16, bold=True, color=COLOR_GREEN)

    nl_flow = [
        ("フルボトル\n空き瓶", COLOR_BLACK),
        ("3Dプリンタ\nで加工", COLOR_GREEN_DARK),
        ("ナイトライト\n完成", COLOR_GREEN),
        ("TikTok動画\nで拡散", COLOR_RED),
        ("販売\n(各チャネル)", COLOR_BLACK),
    ]
    for i, (text, color) in enumerate(nl_flow):
        x = Inches(0.8) + i * Inches(2.4)
        add_rect(slide, x, Inches(5.0), Inches(1.8), Inches(1.0), color,
                 text=text, font_size=11, font_color=COLOR_WHITE, bold=True)
        if i < len(nl_flow) - 1:
            add_arrow(slide, x + Inches(1.85), Inches(5.35), Inches(0.45), Inches(0.25))

    add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.3),
                "香水の「試す → 使う → 残る → 再生 → 広がる」全ライフサイクルをビジネス化",
                font_size=12, bold=True, color=COLOR_GREEN)

    add_footer(slide)
    add_slide_number(slide, 12)


# ══════════════════════════════════════════════════════════════
# スライド13: 月別売上目標（積み上げ棒グラフ）★最重要
# ══════════════════════════════════════════════════════════════
def slide_13_monthly_target():
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "月別売上目標（2026年4月〜2027年1月）", "累計: ¥57,200,000（目標5,000万超え）")

    # データ定義（月, Amazon, TikTok, メルカリ, フルボトル, ナイトライト, 合計）
    months_data = [
        ("4月",  3800, 22,   0,   60,   0,   3882),
        ("5月",  4000, 90,  45,  200,   0,   4335),
        ("6月",  4200, 150,  45,  350,  35,   4780),
        ("7月",  4000, 270, 150,  400,  35,   4855),
        ("8月",  3800, 250, 130,  350,  70,   4600),
        ("9月",  4200, 350, 200,  450,  70,   5270),
        ("10月", 5000, 700, 500,  600, 105,   6905),
        ("11月", 5500, 900, 675,  600, 250,   7925),
        ("12月", 6000, 1200, 900, 800, 348,   9248),
        ("1月",  3800, 600, 450,  400, 150,   5400),
    ]

    # チャネル色
    ch_colors = [COLOR_GREEN, COLOR_BLUE, COLOR_PURPLE, COLOR_RED, COLOR_ORANGE]
    ch_names = ["Amazon", "TikTok", "メルカリ", "フルボトル", "ナイトライト"]

    # 積み上げ棒グラフ
    max_total = 9248  # 12月ピーク
    bar_max_h = Inches(3.8)
    bar_w = Inches(0.85)
    gap = Inches(0.25)
    start_x = Inches(0.5)
    base_y = Inches(5.5)

    for mi, (month, amz, tik, mer, fb, nl, total) in enumerate(months_data):
        x = start_x + mi * (bar_w + gap)
        channels = [amz, tik, mer, fb, nl]
        cum_h = 0

        for ci, val in enumerate(channels):
            if val <= 0:
                continue
            ratio = val / max_total
            seg_h = bar_max_h * ratio
            seg_y = base_y - cum_h - seg_h
            add_rect(slide, x, seg_y, bar_w, seg_h, ch_colors[ci])
            cum_h += seg_h

        # 合計ラベル
        total_h = bar_max_h * (total / max_total)
        add_textbox(slide, x - Inches(0.1), base_y - total_h - Inches(0.3),
                    bar_w + Inches(0.2), Inches(0.25),
                    f"¥{total/1000:.1f}M" if total >= 1000 else f"¥{total}K",
                    font_size=8, bold=True, color=COLOR_BLACK, alignment=PP_ALIGN.CENTER)

        # 月ラベル
        add_textbox(slide, x - Inches(0.05), base_y + Inches(0.05), bar_w + Inches(0.1), Inches(0.25),
                    month, font_size=9, bold=True, color=COLOR_BLACK, alignment=PP_ALIGN.CENTER)

    # 凡例
    legend_y = Inches(5.9)
    for i, (name, color) in enumerate(zip(ch_names, ch_colors)):
        lx = Inches(0.5) + i * Inches(2.3)
        add_rect(slide, lx, legend_y, Inches(0.3), Inches(0.2), color)
        add_textbox(slide, lx + Inches(0.35), legend_y - Inches(0.02), Inches(1.8), Inches(0.25),
                    name, font_size=9, color=COLOR_BODY)

    # 注釈
    notes = [
        ("※7-8月: 夏の需要減を反映", COLOR_GRAY),
        ("※12月: ギフト需要ピーク ¥9.2M", COLOR_RED),
        ("※1月: 年末反動で減少", COLOR_GRAY),
    ]
    for i, (note, color) in enumerate(notes):
        add_textbox(slide, Inches(0.5), Inches(6.3) + i * Inches(0.22), Inches(5), Inches(0.22),
                    note, font_size=8, color=color)

    # 右側サマリー
    add_rect(slide, Inches(8.0), Inches(6.1), Inches(4.8), Inches(0.7), COLOR_GREEN,
             font_size=10, font_color=COLOR_WHITE)
    shape = slide.shapes[-1]
    tf = shape.text_frame
    tf.paragraphs[0].text = "10ヶ月累計: ¥57,200,000"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_WHITE
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_footer(slide)
    add_slide_number(slide, 13)


# ══════════════════════════════════════════════════════════════
# スライド14: チャネル別KPI
# ══════════════════════════════════════════════════════════════
def slide_14_channel_kpi():
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "チャネル別KPI")

    # Amazon KPI
    add_rect(slide, Inches(0.5), Inches(1.3), Inches(3.8), Inches(0.4), COLOR_GREEN,
             text="Amazon", font_size=14, font_color=COLOR_WHITE, bold=True)
    amz_kpis = [
        ("TACOS", "21.9%", "15%"),
        ("月間注文数", "3,500件", "4,500件"),
        ("ROAS", "3.58", "4.5"),
    ]
    for i, (label, current, target) in enumerate(amz_kpis):
        y = Inches(1.85) + i * Inches(0.35)
        add_textbox(slide, Inches(0.6), y, Inches(1.2), Inches(0.3),
                    label, font_size=10, bold=True, color=COLOR_BODY)
        add_textbox(slide, Inches(1.8), y, Inches(0.8), Inches(0.3),
                    current, font_size=10, color=COLOR_BLACK)
        add_textbox(slide, Inches(2.6), y, Inches(0.3), Inches(0.3),
                    "→", font_size=10, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(2.9), y, Inches(0.8), Inches(0.3),
                    target, font_size=10, bold=True, color=COLOR_GREEN)

    # TikTok Shop KPI
    add_rect(slide, Inches(0.5), Inches(3.0), Inches(3.8), Inches(0.4), COLOR_BLUE,
             text="TikTok Shop", font_size=14, font_color=COLOR_WHITE, bold=True)
    tik_kpis = [
        "5月: SPS 30件達成 → プロモーション対象に",
        "7月: レビュー50件達成",
        "10月: 月間1,000件（MELL比50%）",
        "動画: 毎日1本、再生数2,000→5,000回",
    ]
    for i, kpi in enumerate(tik_kpis):
        y = Inches(3.55) + i * Inches(0.3)
        add_textbox(slide, Inches(0.6), y, Inches(3.7), Inches(0.3),
                    kpi, font_size=9, color=COLOR_BODY)

    # メルカリShop KPI
    add_rect(slide, Inches(4.8), Inches(1.3), Inches(3.8), Inches(0.4), COLOR_PURPLE,
             text="メルカリShop", font_size=14, font_color=COLOR_WHITE, bold=True)
    mer_kpis = [
        "6月: レビュー100件",
        "8月: フォロワー500人",
        "10月: 日販20件（月600件）、レビュー500件",
    ]
    for i, kpi in enumerate(mer_kpis):
        y = Inches(1.85) + i * Inches(0.3)
        add_textbox(slide, Inches(4.9), y, Inches(3.7), Inches(0.3),
                    kpi, font_size=10, color=COLOR_BODY)

    # LINE → フルボトル KPI
    add_rect(slide, Inches(4.8), Inches(3.0), Inches(3.8), Inches(0.4), COLOR_RED,
             text="LINE → フルボトル", font_size=14, font_color=COLOR_WHITE, bold=True)
    line_kpis = [
        ("登録率", "5.3%", "8%"),
        ("シナリオ開封率", "-", "60%"),
        ("フルボトル転換率", "-", "20%"),
    ]
    for i, (label, current, target) in enumerate(line_kpis):
        y = Inches(3.55) + i * Inches(0.3)
        add_textbox(slide, Inches(4.9), y, Inches(1.5), Inches(0.3),
                    label, font_size=10, bold=True, color=COLOR_BODY)
        add_textbox(slide, Inches(6.4), y, Inches(0.6), Inches(0.3),
                    current, font_size=10, color=COLOR_BLACK)
        add_textbox(slide, Inches(7.0), y, Inches(0.3), Inches(0.3),
                    "→", font_size=10, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(7.3), y, Inches(0.8), Inches(0.3),
                    target, font_size=10, bold=True, color=COLOR_GREEN)

    # ナイトライト KPI
    add_rect(slide, Inches(9.1), Inches(1.3), Inches(3.8), Inches(0.4), COLOR_ORANGE,
             text="ナイトライト", font_size=14, font_color=COLOR_WHITE, bold=True)
    nl_kpis = [
        "6月: 製作開始（自社ボトル月50個）",
        "8月: 返送加工サービス開始",
        "11月: 月100個（50個自社+50個返送加工）",
    ]
    for i, kpi in enumerate(nl_kpis):
        y = Inches(1.85) + i * Inches(0.3)
        add_textbox(slide, Inches(9.2), y, Inches(3.7), Inches(0.3),
                    kpi, font_size=10, color=COLOR_BODY)

    # 下部まとめ
    add_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.5), COLOR_GREEN_LIGHT,
             font_size=10, font_color=COLOR_BLACK)
    shape = slide.shapes[-1]
    tf = shape.text_frame
    tf.paragraphs[0].text = "KPIまとめ: 最優先はレビュー獲得"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_BLACK
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    priorities = [
        "1. TikTok SPS 30件（5月） → レビュー50件（7月） → 月間1,000件（10月）",
        "2. メルカリ レビュー100件（6月） → 日販20件（10月）",
        "3. Amazon TACOS 21.9% → 15%（広告効率改善）",
        "4. LINE フルボトル転換率 20%（CRM育成）",
    ]
    for pri in priorities:
        p = tf.add_paragraph()
        p.text = pri
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_BODY
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.LEFT

    add_footer(slide)
    add_slide_number(slide, 14)


# ══════════════════════════════════════════════════════════════
# スライド15: 利益シミュレーション
# ══════════════════════════════════════════════════════════════
def slide_15_profit_simulation():
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "利益シミュレーション", "10ヶ月累計（2026年4月〜2027年1月）")

    # 左: P/Lサマリー
    add_textbox(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.3),
                "10ヶ月累計", font_size=18, bold=True, color=COLOR_GREEN)

    pl_items = [
        ("売上", "¥57,200,000", COLOR_BLACK, True),
        ("粗利", "¥16,223,380", COLOR_GREEN, True),
        ("人件費", "-¥1,600,000", COLOR_RED, False),
        ("", "(パート月16万x10ヶ月)", COLOR_GRAY, False),
        ("営業利益", "¥14,623,380", COLOR_GREEN, True),
        ("営業利益率", "25.6%", COLOR_GREEN, True),
    ]
    for i, (label, value, color, bold) in enumerate(pl_items):
        y = Inches(1.9) + i * Inches(0.5)
        if label == "営業利益":
            sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         Inches(0.8), y - Inches(0.08), Inches(4.5), Inches(0.03))
            sep.fill.solid()
            sep.fill.fore_color.rgb = COLOR_GREEN
            sep.line.fill.background()
        add_textbox(slide, Inches(0.8), y, Inches(2.0), Inches(0.3),
                    label, font_size=14, bold=bold, color=color)
        add_textbox(slide, Inches(3.0), y, Inches(2.5), Inches(0.3),
                    value, font_size=14 if bold else 11, bold=bold, color=color,
                    alignment=PP_ALIGN.RIGHT)

    # 月平均カード
    add_rect(slide, Inches(0.8), Inches(5.0), Inches(2.5), Inches(1.0), COLOR_GREEN,
             font_size=10, font_color=COLOR_WHITE)
    shape = slide.shapes[-1]
    tf = shape.text_frame
    tf.paragraphs[0].text = "月平均売上"
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = RGBColor(0xCC, 0xFF, 0xCC)
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "¥5,720,000"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_WHITE
    p2.font.name = FONT_NAME
    p2.alignment = PP_ALIGN.CENTER

    add_rect(slide, Inches(3.5), Inches(5.0), Inches(2.5), Inches(1.0), COLOR_GREEN_DARK,
             font_size=10, font_color=COLOR_WHITE)
    shape = slide.shapes[-1]
    tf = shape.text_frame
    tf.paragraphs[0].text = "月平均利益"
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = RGBColor(0xCC, 0xFF, 0xCC)
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "¥1,462,338"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_WHITE
    p2.font.name = FONT_NAME
    p2.alignment = PP_ALIGN.CENTER

    # 右: 将来展望
    rx = Inches(7.0)
    add_textbox(slide, rx, Inches(1.4), Inches(5), Inches(0.3),
                "将来展望", font_size=18, bold=True, color=COLOR_GREEN)

    future_items = [
        ("問屋を介さず海外から直接仕入れ", "原価率さらに改善", COLOR_GREEN),
        ("守殿さん復帰（11月〜）", "オペレーション体制強化", COLOR_BLACK),
        ("TikTok Shop成長", "広告費不要の販売チャネル拡大", COLOR_BLUE),
        ("ナイトライト事業", "超高利益率の成長ドライバー", COLOR_ORANGE),
        ("LINE CRM成熟", "リピーター基盤の確立", COLOR_RED),
    ]
    for i, (title, desc, color) in enumerate(future_items):
        y = Inches(1.9) + i * Inches(0.8)
        add_rect(slide, rx, y, Inches(0.15), Inches(0.5), color)
        add_textbox(slide, rx + Inches(0.3), y, Inches(5), Inches(0.3),
                    title, font_size=14, bold=True, color=COLOR_BLACK)
        add_textbox(slide, rx + Inches(0.3), y + Inches(0.35), Inches(5), Inches(0.3),
                    f"→ {desc}", font_size=12, color=COLOR_GRAY)

    # 売上構成の横棒（下部）
    add_textbox(slide, rx, Inches(5.1), Inches(5), Inches(0.3),
                "売上構成比（10ヶ月累計）", font_size=13, bold=True, color=COLOR_GREEN)

    bar_data = [
        ("Amazon", 44200, COLOR_GREEN),
        ("TikTok", 4532, COLOR_BLUE),
        ("メルカリ", 3095, COLOR_PURPLE),
        ("フルボトル", 4210, COLOR_RED),
        ("ナイトライト", 1163, COLOR_ORANGE),
    ]
    bar_max_w = Inches(4.5)
    total_sales = sum(v for _, v, _ in bar_data)
    cum_x = rx
    for name, val, color in bar_data:
        w = bar_max_w * (val / total_sales)
        if w < Inches(0.1):
            w = Inches(0.1)
        add_rect(slide, cum_x, Inches(5.45), w, Inches(0.4), color)
        cum_x += w

    # 凡例
    for i, (name, val, color) in enumerate(bar_data):
        lx = rx + (i % 3) * Inches(1.8)
        ly = Inches(5.95) + (i // 3) * Inches(0.25)
        add_rect(slide, lx, ly, Inches(0.2), Inches(0.15), color)
        pct = val / total_sales * 100
        add_textbox(slide, lx + Inches(0.25), ly - Inches(0.03), Inches(1.5), Inches(0.2),
                    f"{name} {pct:.0f}%", font_size=8, color=COLOR_BODY)

    add_footer(slide)
    add_slide_number(slide, 15)


# ══════════════════════════════════════════════════════════════
# スライド16: ロードマップ
# ══════════════════════════════════════════════════════════════
def slide_16_roadmap():
    slide = prs.slides.add_slide(blank_layout)
    add_section_title(slide, "ロードマップ", "2026年4月〜2027年1月")

    # タイムラインカラム
    phases = [
        ("2026年4月\n（現在）", [
            ("✅", "Amazon月間3,500件・1位"),
            ("✅", "自社サイトリニューアル完了"),
            ("✅", "AI香水診断稼働"),
            ("✅", "LINE CRMシナリオ設計完了"),
            ("★", "LINE配信 明日ローンチ"),
            ("🔄", "TikTok Shop 25件/月"),
        ], COLOR_GREEN),
        ("2026年\n5〜6月", [
            ("□", "TikTok SPS30件→プロモ解放"),
            ("□", "メルカリShop本格稼働"),
            ("□", "Threads/Instagram開始"),
            ("□", "ナイトライト製作開始"),
            ("□", "ブログ100本"),
        ], COLOR_GREEN_DARK),
        ("2026年\n7〜9月", [
            ("□", "TikTokレビュー50件"),
            ("□", "メルカリレビュー100件"),
            ("□", "ナイトライト返送加工開始"),
            ("□", "フルボトル月間¥40-50万"),
        ], COLOR_BLACK),
        ("2026年\n10〜12月", [
            ("□", "TikTok月間1,000件"),
            ("□", "メルカリ月間750件"),
            ("□", "12月ピーク: 月商¥925万"),
            ("□", "ナイトライト月100個"),
            ("□", "11月: 守殿さん復帰"),
        ], COLOR_RED),
        ("2027年\n1月", [
            ("□", "累計売上¥5,720万達成"),
            ("□", "海外直接仕入れの検討開始"),
        ], COLOR_BLACK),
    ]

    col_w = Inches(2.3)
    gap = Inches(0.15)
    start_x = Inches(0.5)

    for pi, (period, items, color) in enumerate(phases):
        x = start_x + pi * (col_w + gap)

        # 期間ヘッダー
        add_rect(slide, x, Inches(1.3), col_w, Inches(0.6), color,
                 text=period, font_size=12, font_color=COLOR_WHITE, bold=True)

        # 項目
        for ii, (icon, text) in enumerate(items):
            y = Inches(2.1) + ii * Inches(0.42)
            icon_color = COLOR_GREEN if icon == "✅" else (COLOR_RED if icon == "★" else (COLOR_ORANGE if icon == "🔄" else COLOR_GRAY))

            if icon in ("✅", "★"):
                bg = COLOR_GREEN_LIGHT
            else:
                bg = COLOR_WHITE

            shape = add_rect(slide, x, y, col_w, Inches(0.38), bg,
                             font_size=9, font_color=COLOR_BLACK)
            tf = shape.text_frame
            tf.paragraphs[0].text = f"{icon} {text}"
            tf.paragraphs[0].font.size = Pt(9)
            tf.paragraphs[0].font.color.rgb = COLOR_GREEN_DARK if icon in ("✅", "★") else COLOR_BODY
            tf.paragraphs[0].font.bold = icon == "★"
            tf.paragraphs[0].font.name = FONT_NAME
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT

        # 矢印（最後以外）
        if pi < len(phases) - 1:
            add_arrow(slide, x + col_w, Inches(1.5), gap, Inches(0.2), COLOR_GRAY)

    # 下部まとめ
    add_rect(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.2), COLOR_GREEN,
             font_size=10, font_color=COLOR_WHITE)
    shape = slide.shapes[-1]
    tf = shape.text_frame
    tf.paragraphs[0].text = "10ヶ月で売上¥5,720万・営業利益¥1,462万を目指す"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_WHITE
    tf.paragraphs[0].font.name = FONT_NAME
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Amazon基盤 × TikTok/メルカリ拡大 × LINE CRM育成 × ナイトライト新規事業"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(0xCC, 0xFF, 0xCC)
    p2.font.name = FONT_NAME
    p2.alignment = PP_ALIGN.CENTER


    add_footer(slide)
    add_slide_number(slide, 16)


# ══════════════════════════════════════════════════════════════
# 実行
# ══════════════════════════════════════════════════════════════
slide_01_cover()
slide_02_overview()
slide_03_business_model()
slide_04_competitor()
slide_05_amazon_sales()
slide_06_unit_economics()
slide_07_monthly_pl()
slide_08_funnel()
slide_09_customer()
slide_10_channels()
slide_11_line_crm()
slide_12_site_sns()
slide_13_monthly_target()
slide_14_channel_kpi()
slide_15_profit_simulation()
slide_16_roadmap()

output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "COLLEGRANCE_事業全体像.pptx")
prs.save(output_path)
print(f"✅ 生成完了: {output_path}")
print(f"   スライド数: {len(prs.slides)}")
